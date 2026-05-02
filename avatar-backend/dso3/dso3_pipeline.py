"""
DSO3 Pipeline - Compliance Indication Rewriter
Serializable pipeline for joblib export
"""

import json
import os
import re
from datetime import datetime
from typing import Dict, Optional
import pandas as pd


class DSO3Config:
    """Configuration for DSO3 Engine"""
    MODEL = "llama-3.3-70b-versatile"
    TEMPERATURE = 0.1
    MAX_TOKENS = 2048


def get_compliance_prompt(product_data: Dict) -> str:
    """Build the RAG prompt with French product data"""
    return f"""Vous êtes un expert en conformité réglementaire pour les compléments alimentaires et les produits cosmétiques (directives FDA/EMA et réglementation française).

Voici les données du produit:

NOM DU PRODUIT: {product_data.get('produit', 'Inconnu')}
GAMME: {product_data.get('gamme', 'Non spécifié')}
INDICATION ORIGINALE: {product_data.get('indication', '')}
COMPOSITION: {product_data.get('composition', 'Non spécifiée')}
POSOLOGIE: {product_data.get('posologie', "Suivre les instructions sur l'emballage")}

---

TÂCHE 1: DÉTECTER LES PHRASES RISQUÉES
Analysez l'INDICATION ORIGINALE et identifiez:
- Verbes interdits (traite, guérit, prévient, élimine, soigne, combat, détruit)
- Allégations de maladie (acné, eczéma, infection, inflammation, douleur, mycose, constipation, anémie, stress, anxiété, migraine, insomnie)
- Superlatifs/absolus (meilleur, parfait, 100%, garanti, tous, jamais)
- Absence de qualificatifs (devrait avoir "peut aider", "contribue à", "soutient", "maintient")

TÂCHE 2: RÉÉCRIRE EN INDICATION CONFORME
Règles:
- Remplacer les allégations de maladie par des descriptions de fonction corporelle
- Ajouter des qualificatifs: "peut aider", "contribue à", "soutient", "favorise", "maintient"
- Supprimer toutes les allégations absolues
- Pour les compléments alimentaires: aucune mention de maladie
- Pour les cosmétiques: aucune allégation de traitement de maladie

TÂCHE 3: AJOUTER LES MENTIONS LÉGALES OBLIGATOIRES

Pour les COMPLÉMENTS ALIMENTAIRES (produits à avaler):
"⚠️ Ce produit ne peut pas diagnostiquer, traiter, guérir ou prévenir une quelconque maladie. Consultez votre médecin avant utilisation si vous êtes enceinte, allaitez, prenez des médicaments ou avez un problème médical. Ne pas dépasser la dose recommandée."

Pour les PRODUITS COSMÉTIQUES (application sur la peau):
"⚠️ Usage externe uniquement. En cas d'irritation, arrêtez l'utilisation. Évitez le contact avec les yeux. Ce produit n'est pas destiné à traiter des maladies de la peau."

Pour les PRODUITS ANTI-MOUSTIQUES/INSECTES:
"⚠️ Usage externe uniquement. Ne pas appliquer sur les muqueuses ou les plaies. Tenir hors de portée des enfants."

TÂCHE 4: CRÉER UNE PRÉSENTATION DE 7 DIAPOSITIVES
Générez ces diapositives exactement en FRANÇAIS:

1. TITRE: [Nom du produit] — [Catégorie de bénéfice]
2. INTRODUCTION: À qui s'adresse ce produit
3. BÉNÉFICE PRINCIPAL: L'indication conforme
4. COMMENT ÇA MARCHE: Basé sur la composition
5. MODE D'EMPLOI: La posologie
6. SÉCURITÉ: Les mentions légales
7. CONCLUSION: Avertissement légal complet

---

RETOURNEZ UNIQUEMENT DU JSON VALIDE avec cette structure:

{{
  "product_name": "string",
  "risk_detection": {{
    "forbidden_verbs": [],
    "disease_claims": [],
    "superlatives": [],
    "missing_qualifiers": false
  }},
  "compliant_indication": "string",
  "full_compliant_statement": "string",
  "presentation": {{
    "slide1_title": "string",
    "slide2_introduction": "string",
    "slide3_key_benefit": "string",
    "slide4_how_it_works": "string",
    "slide5_usage": "string",
    "slide6_safety": "string",
    "slide7_closing": "string"
  }}
}}"""


class DSO3Engine:
    """
    Main DSO3 Engine — serializable via joblib.
    All Groq calls are deferred to call-time (api_key injected at runtime).
    """

    def __init__(self):
        self.model = DSO3Config.MODEL
        self.temperature = DSO3Config.TEMPERATURE
        self.max_tokens = DSO3Config.MAX_TOKENS
        # Config is stored; client is built at runtime to avoid serialization issues
        self._config = {
            "model": self.model,
            "temperature": self.temperature,
            "max_tokens": self.max_tokens,
        }

    def _get_client(self, api_key: str):
        from groq import Groq
        return Groq(api_key=api_key)

    def process_product(self, product_data: Dict, api_key: str) -> Dict:
        """Call Groq and return compliance JSON"""
        if not product_data.get("indication"):
            return {
                "error": "Aucune indication fournie",
                "product_name": product_data.get("produit", "Inconnu"),
            }

        client = self._get_client(api_key)
        prompt = get_compliance_prompt(product_data)

        try:
            response = client.chat.completions.create(
                model=self.model,
                messages=[
                    {
                        "role": "system",
                        "content": "Vous êtes un expert en conformité réglementaire. Retournez uniquement du JSON valide.",
                    },
                    {"role": "user", "content": prompt},
                ],
                temperature=self.temperature,
                max_tokens=self.max_tokens,
                response_format={"type": "json_object"},
            )
            result = json.loads(response.choices[0].message.content)
            result["processed_at"] = datetime.now().isoformat()
            result["model_used"] = self.model
            return result
        except Exception as e:
            return {
                "error": str(e),
                "product_name": product_data.get("produit", "Inconnu"),
                "processed_at": datetime.now().isoformat(),
            }

    def generate_pptx_bytes(self, product_data: Dict, result: Dict) -> Optional[bytes]:
        """Generate PowerPoint and return as bytes (for FastAPI response)"""
        try:
            import io
            from pptx import Presentation

            prs = Presentation()
            title_layout = prs.slide_layouts[0]
            content_layout = prs.slide_layouts[1]
            pdata = result.get("presentation", {})
            name = product_data.get("produit", "Inconnu")

            slides_cfg = [
                (title_layout, name, "Présentation du produit conforme", is_title := True),
                (content_layout, "Introduction", pdata.get("slide2_introduction", ""), False),
                (content_layout, "Bénéfice Principal", pdata.get("slide3_key_benefit", ""), False),
                (content_layout, "Comment ça marche", pdata.get("slide4_how_it_works", ""), False),
                (content_layout, "Mode d'emploi", pdata.get("slide5_usage", ""), False),
                (content_layout, "Informations de sécurité", pdata.get("slide6_safety", ""), False),
                (content_layout, "Mentions légales importantes", pdata.get("slide7_closing", ""), False),
            ]

            for layout, title, body, _ in slides_cfg:
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = title
                if slide.placeholders and len(slide.placeholders) > 1:
                    slide.placeholders[1].text = body

            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)
            return buf.read()
        except Exception as e:
            print(f"PowerPoint error: {e}")
            return None

    def generate_compliance(self, product_data: Dict, api_key: str) -> Dict:
        """
        Full pipeline: process + optional pptx bytes.
        Returns a dict ready for FastAPI JSON response.
        """
        result = self.process_product(product_data, api_key)

        if result.get("error"):
            return {"success": False, "error": result["error"], "product_name": product_data.get("produit", "Inconnu")}

        pptx_bytes = self.generate_pptx_bytes(product_data, result)

        return {
            "success": True,
            "product_name": product_data.get("produit", "Inconnu"),
            "compliant_indication": result.get("compliant_indication", ""),
            "full_compliant_statement": result.get("full_compliant_statement", ""),
            "risk_detection": result.get("risk_detection", {}),
            "presentation_slides": result.get("presentation", {}),
            "pptx_available": pptx_bytes is not None,
            "processed_at": result.get("processed_at", ""),
            "model_used": result.get("model_used", ""),
            # Internal only — not serialized to JSON
            "_pptx_bytes": pptx_bytes,
        }
